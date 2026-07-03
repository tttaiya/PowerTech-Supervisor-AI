from __future__ import annotations

import logging
import tempfile
import time
from pathlib import Path

logger = logging.getLogger(__name__)

from app.schemas import ParsedBlock, TaskPayload
from app.services.ocr import ocr_image
from app.services.path_guard import MAX_PDF_PAGES
from .text_utils import (
    clean_document_text,
    clean_pdf_pages,
    clean_text,
    ensure_file_exists,
    guess_extension,
    read_text_file,
    update_chapter_path,
)

IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "bmp", "webp"}
TEXT_EXTENSIONS = {"txt", "md"}
P0_EXTENSIONS = {"pdf", "docx", *TEXT_EXTENSIONS, *IMAGE_EXTENSIONS}
P1_EXTENSIONS = {"pptx", "xlsx"}
SUPPORTED_EXTENSIONS = P0_EXTENSIONS | P1_EXTENSIONS


def parse_document(file_path: str, extension: str | None, payload: TaskPayload) -> tuple[str, list[ParsedBlock], str]:
    """统一解析入口。返回 parsed_text、blocks、extension。"""
    path = ensure_file_exists(file_path)
    ext = guess_extension(str(path), extension)
    parse_backend = (payload.parse_backend or "pymupdf").lower()

    # Tika 本期不做真实实现，只保留 parseBackend 字段，避免组长后面扩展时改接口。
    if parse_backend not in {"pymupdf", "self", "auto", "paddleocr"}:
        raise NotImplementedError(
            f"parseBackend={payload.parse_backend} 已预留，但本期最小版本只实现 pymupdf/self/auto/paddleocr。"
        )

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"暂不支持的文件类型：{ext}，当前支持：{sorted(SUPPORTED_EXTENSIONS)}")

    if ext == "pdf":
        blocks = parse_pdf(path, payload)
    elif ext == "docx":
        blocks = parse_docx(path)
    elif ext in TEXT_EXTENSIONS:
        blocks = parse_text(path)
    elif ext in IMAGE_EXTENSIONS:
        blocks = parse_image(path, payload)
    elif ext == "pptx":
        blocks = parse_pptx(path)
    elif ext == "xlsx":
        blocks = parse_xlsx(path)
    else:
        raise ValueError(f"未知文件类型：{ext}")

    blocks = [b for b in blocks if b.content and b.content.strip()]
    if not blocks:
        raise ValueError("解析结果为空，请检查文件内容或 OCR 配置。")

    parsed_text = "\n\n".join(block.content for block in blocks if block.content.strip())
    return parsed_text, blocks, ext


def parse_pdf(path: Path, payload: TaskPayload) -> list[ParsedBlock]:
    """PDF：先提取文本层，只对无文本页面执行 OCR。"""
    try:
        import fitz  # PyMuPDF
    except Exception as exc:
        raise RuntimeError("缺少 PyMuPDF，请执行：pip install PyMuPDF") from exc

    try:
        doc = fitz.open(str(path))
    except Exception as exc:
        raise ValueError("PDF 文件损坏、加密或无法打开") from exc

    with doc:
        page_limit = min(payload.max_pdf_pages, MAX_PDF_PAGES)

        if doc.page_count > page_limit:
            raise ValueError(
                f"PDF 页数超过限制：{doc.page_count} 页（最大 {page_limit} 页）"
            )

        if doc.page_count == 0:
            raise ValueError("PDF 不包含可解析页面")

        logger.info(
            "PDF parse start path=%s pages=%d enableOcr=%s",
            path,
            doc.page_count,
            payload.enable_ocr,
        )

        # 第一步：提取原始文本层。
        raw_page_texts: list[str] = []
        extract_started = time.time()

        for page_index, page in enumerate(doc, start=1):
            page_started = time.time()
            text = page.get_text("text") or ""
            raw_page_texts.append(text)

            if (
                page_index == 1
                or page_index % 10 == 0
                or page_index == doc.page_count
            ):
                logger.info(
                    "PDF text progress page=%d/%d chars=%d pageMs=%d totalMs=%d",
                    page_index,
                    doc.page_count,
                    len(text),
                    int((time.time() - page_started) * 1000),
                    int((time.time() - extract_started) * 1000),
                )

        page_texts = clean_pdf_pages(raw_page_texts)
        total_chars = sum(len(text or "") for text in page_texts)

        logger.info(
            "PDF text extraction finished pages=%d totalChars=%d",
            doc.page_count,
            total_chars,
        )

        # 第二步：仅 OCR 没有有效文本的页面。
        if payload.enable_ocr:
            ocr_page_indexes = [
                index
                for index, text in enumerate(page_texts)
                if len((text or "").strip()) < 10
            ]

            logger.info(
                "PDF OCR fallback pages=%d/%d",
                len(ocr_page_indexes),
                doc.page_count,
            )

            if ocr_page_indexes:
                with tempfile.TemporaryDirectory(
                    prefix="km_pdf_ocr_"
                ) as tmp_dir:
                    ocr_started = time.time()

                    for order, page_index in enumerate(
                        ocr_page_indexes,
                        start=1,
                    ):
                        page_started = time.time()
                        page = doc.load_page(page_index)

                        # 原来是 2 倍，CPU OCR 较慢，调整为 1.5 倍。
                        pix = page.get_pixmap(
                            matrix=fitz.Matrix(1.5, 1.5),
                            alpha=False,
                        )

                        image_path = (
                            Path(tmp_dir) / f"page_{page_index + 1}.png"
                        )
                        pix.save(str(image_path))

                        text, _ = ocr_image(image_path)
                        text = text or ""

                        if text.strip():
                            page_texts[page_index] = text

                        logger.info(
                            "PDF OCR progress page=%d/%d "
                            "ocrIndex=%d/%d chars=%d "
                            "pageMs=%d totalMs=%d",
                            page_index + 1,
                            doc.page_count,
                            order,
                            len(ocr_page_indexes),
                            len(text),
                            int((time.time() - page_started) * 1000),
                            int((time.time() - ocr_started) * 1000),
                        )

                page_texts = clean_pdf_pages(page_texts)

        blocks: list[ParsedBlock] = []

        for page_index, text in enumerate(page_texts, start=1):
            text = text or ""

            if not text.strip():
                continue

            has_original_text = len(
                (raw_page_texts[page_index - 1] or "").strip()
            ) >= 10

            blocks.append(
                ParsedBlock(
                    content=text,
                    pageNo=page_index,
                    blockType=(
                        "paragraph" if has_original_text else "ocr"
                    ),
                    chapterPath=None,
                    charCount=len(text),
                )
            )

        logger.info(
            "PDF parse finished pages=%d blocks=%d chars=%d",
            doc.page_count,
            len(blocks),
            sum(len(block.content) for block in blocks),
        )

        return blocks


def parse_docx(path: Path) -> list[ParsedBlock]:
    try:
        from docx import Document
    except Exception as exc:
        raise RuntimeError("缺少 python-docx，请执行：pip install python-docx") from exc

    doc = Document(str(path))
    blocks: list[ParsedBlock] = []
    heading_levels: list[str] = []

    for paragraph in doc.paragraphs:
        text = clean_document_text(paragraph.text or "")
        if not text:
            continue

        # 优先使用 Word 自带标题样式判断章节层级。
        style_name = (paragraph.style.name or "") if paragraph.style else ""
        if style_name.lower().startswith("heading") or style_name.startswith("标题"):
            # 尝试取标题级别
            level = 1
            for token in style_name.split():
                if token.isdigit():
                    level = int(token)
                    break
            idx = max(0, level - 1)
            while len(heading_levels) <= idx:
                heading_levels.append("")
            heading_levels[idx] = text
            del heading_levels[idx + 1 :]
        else:
            # 没有标题样式时，用“第一章 / 一、 / 1.1”这类文本规则兜底。
            update_chapter_path(text, heading_levels)

        chapter_path = "/".join([item for item in heading_levels if item]) or None
        blocks.append(
            ParsedBlock(
                content=text,
                pageNo=None,
                blockType="paragraph",
                chapterPath=chapter_path,
                charCount=len(text),
            )
        )

    # 表格也做简单提取，便于后续检索。
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [clean_text(cell.text or "") for cell in row.cells]
            rows.append("\t".join(cells))
        table_text = clean_text("\n".join(rows))
        if table_text:
            blocks.append(
                ParsedBlock(
                    content=table_text,
                    pageNo=None,
                    blockType="table",
                    chapterPath="/".join([item for item in heading_levels if item]) or None,
                    charCount=len(table_text),
                )
            )
    return blocks


def parse_text(path: Path) -> list[ParsedBlock]:
    text = clean_document_text(read_text_file(path))
    return [ParsedBlock(content=text, pageNo=1, blockType="paragraph", chapterPath=None, charCount=len(text))] if text else []


def parse_image(path: Path, payload: TaskPayload) -> list[ParsedBlock]:
    if not payload.enable_ocr:
        raise RuntimeError("图片文件必须启用 OCR 才能解析。")
    text, _ = ocr_image(path)
    text = clean_document_text(text)
    return [ParsedBlock(content=text, pageNo=1, blockType="ocr", chapterPath=None, charCount=len(text))] if text else []


def parse_pptx(path: Path) -> list[ParsedBlock]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError("缺少 python-pptx，请执行：pip install python-pptx") from exc

    prs = Presentation(str(path))
    blocks: list[ParsedBlock] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        text = clean_document_text("\n".join(texts))
        if text:
            blocks.append(ParsedBlock(content=text, pageNo=index, blockType="paragraph", charCount=len(text)))
    return blocks


def parse_xlsx(path: Path) -> list[ParsedBlock]:
    try:
        import openpyxl
    except Exception as exc:
        raise RuntimeError("缺少 openpyxl，请执行：pip install openpyxl") from exc

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    blocks: list[ParsedBlock] = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) if cell is not None else "" for cell in row]
            if any(value.strip() for value in values):
                rows.append("\t".join(values))
        text = clean_text(f"工作表：{sheet.title}\n" + "\n".join(rows))
        if text:
            blocks.append(ParsedBlock(content=text, pageNo=None, blockType="table", chapterPath=sheet.title, charCount=len(text)))
    return blocks
